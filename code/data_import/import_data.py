from qdrant_client import QdrantClient
from qdrant_client.http.models import Distance, VectorParams
from qdrant_client.http.models import PointStruct
from qdrant_client.models import Filter, FieldCondition
import tqdm
import PyPDF2
import uuid
import docx
import sys
import re
import os
from text2vec import SentenceModel
from PyPDF2.errors import PdfReadError
import time
import json
import argparse
from pptx import Presentation

"""
找到重复的文件，并且移动到指定的目录中
"""
import hashlib
import shutil

# 计算文件的MD5哈希值。
def file_hash(file_path):
    with open(file_path, 'rb') as f:
        file_data = f.read()
        file_hash = hashlib.md5(file_data).hexdigest()
    return file_hash

# 在源文件夹中查找重复文件并将其移动到指定的重复文件夹。
def find_duplicate_files_and_move(src_folder, duplicate_folder):
    if not os.path.exists(duplicate_folder):
        os.makedirs(duplicate_folder)

    file_hashes = {}
    for root, _, files in os.walk(src_folder):
        for file in files:
            file_path = os.path.join(root, file)
            current_file_hash = file_hash(file_path)

            if current_file_hash in file_hashes:
                existing_file_path = file_hashes[current_file_hash]
                if len(file) > len(os.path.basename(existing_file_path)):
                    duplicate_file_path = os.path.join(duplicate_folder, file)
                    shutil.move(file_path, duplicate_file_path)
                    print(f"Moved duplicate file: {file_path} to {duplicate_file_path}")
                else:
                    duplicate_file_path = os.path.join(duplicate_folder, os.path.basename(existing_file_path))
                    shutil.move(existing_file_path, duplicate_file_path)
                    print(f"Moved duplicate file: {existing_file_path} to {duplicate_file_path}")
                    file_hashes[current_file_hash] = file_path
            else:
                file_hashes[current_file_hash] = file_path


# 过滤文本中的中文字符、标点符号和空格。
def filter_chinese_and_punctuations(text):
    # 定义正则表达式，匹配中文、数字和标准的标点符号、英文字符和空格、回车符、制表符等，以及键盘上数字哪一行的所有符号
    pattern = re.compile(r'[\u4e00-\u9fa5\d ，。！？、；：‘’“”（）《》【】\[\]【】a-zA-Z,.!?;:\'"/\\\{\}\(\)\<\>\+\-\*/=~=^`|&#%@_\n\r\t]+')
    # 使用正则表达式过滤文本
    result = pattern.findall(text)
    filtered_text = ''.join(result)
    # 汉字中间的空格给去掉，pdf转译出来有很多空格
    filtered_text = re.sub(r'([\u4e00-\u9fa5])\s+([\u4e00-\u9fa5])', r'\1\2', filtered_text)
    return filtered_text

os.environ["PROTOCOL_BUFFERS_PYTHON_IMPLEMENTATION"] = "python"

# 使用预训练模型生成文本摘要。
from transformers import AutoTokenizer, AutoModelForSeq2SeqLM
WHITESPACE_HANDLER = lambda k: re.sub('\s+', ' ', re.sub('\n+', ' ', k.strip()))
summary_model_name = "csebuetnlp/mT5_multilingual_XLSum"
tokenizer = AutoTokenizer.from_pretrained(summary_model_name)
summary_model = AutoModelForSeq2SeqLM.from_pretrained(summary_model_name)
def get_summary(text, summary_length=200):
    input_ids = tokenizer(
        [WHITESPACE_HANDLER(text)],
        return_tensors="pt",
        padding="max_length",
        truncation=True,
        max_length=512
    )["input_ids"]

    output_ids = summary_model.generate(
        input_ids=input_ids,
        max_length=summary_length, #84
        no_repeat_ngram_size=2,
        num_beams=4
    )[0]

    summary = tokenizer.decode(
        output_ids,
        skip_special_tokens=True,
        clean_up_tokenization_spaces=False
    )
    return summary

# 使用预训练模型将给定文本转换为句向量表示
namespace = uuid.NAMESPACE_URL
embed_model = SentenceModel('shibing624/text2vec-base-chinese')
def to_embeddings(text):
    embeddings = embed_model.encode([text])
    return embeddings[0].tolist()

# 将文本分割为最多包含1000个token的块，用来做Qdrant的索引。
from langchain.text_splitter import TokenTextSplitter
def split_text_to_chunks(text, max_tokens=1000):
    text = filter_chinese_and_punctuations(text)
    text = text.replace("\n\n","\n").replace("  "," ").replace("\t\t","\t").replace("..",".")
    text_splitter = TokenTextSplitter(chunk_size=max_tokens, chunk_overlap=0)
    texts = text_splitter.split_text(text)
    return texts

# PDF Reader

def get_pdf_chunks(file_path, max_chunk_size=1000):
    chunks = []
    try:
        with open(file_path, "rb") as f:
            pdf_reader = PyPDF2.PdfReader(f)
            num_pages = len(pdf_reader.pages)
            pages = [pdf_reader.pages[i].extract_text() for i in range(num_pages)]

        full_text = "".join(pages)
        chunks = split_text_to_chunks(full_text, max_chunk_size)
    except PdfReadError:
        print(f"Warning: Unable to read the PDF file '{file_path}', EOF marker not found. Skipping this file.")
    return chunks

# pptx Reader

def read_pptx_text(pptx_file):
    presentation = Presentation(pptx_file)
    text = []
    try:
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text.append(paragraph.text)
    except Exception as e:
        print(f"Warning: Unable to read the PDF file '{pptx_file}', Exception found. Skipping this file.")
        print(e)

    return ' '.join(text)

# 更新状态文件以记录处理过程。
def update_processing_status_file(status_file_path, file_path, last_processed_chunk, is_completed):
    with open(status_file_path, 'r') as f:
        status_data = json.load(f)
    
    status_data[file_path] = {
        'last_processed_chunk': last_processed_chunk,
        'is_completed': is_completed
    }

    with open(status_file_path, 'w') as f:
        json.dump(status_data, f)
        
# 从状态文件中获取给定文件的处理状态。
def get_file_processing_status(status_file_path, file_path):
    with open(status_file_path, 'r') as f:
        status_data = json.load(f)
    return status_data.get(file_path, None)
        
"""
主程序，一直运行，自动索引新的文件
main_loop函数的处理流程：

1. 扫描给定的基础目录，找到所有支持的文件类型（.pdf, .docx, .pptx）。

2. 对于每个找到的文件：
   a. 检查状态文件，获取文件的处理状态。
   b. 如果文件已完成处理，跳过此文件。
   c. 如果文件尚未完成处理，从上次处理的位置开始继续处理。

3. 根据文件类型分别提取文本内容：
   a. 对于PDF文件，使用get_pdf_chunks函数提取文本并将其分割为子字符串。
   b. 对于DOCX文件，使用python-docx库提取文本。
   c. 对于PPTX文件，使用read_pptx_text函数提取文本。

4. 对提取到的文本进行预处理：
   a. 使用filter_chinese_and_punctuations函数过滤掉中文字符、标点符号和空格。
   b. 使用split_text_to_chunks函数将文本分割为最多包含1000个字符的子字符串。

5. 对每个文本子字符串执行以下操作：
   a. 使用预训练模型生成摘要。
   b. 将摘要转换为句向量表示。
   c. 将摘要及其句向量表示添加到指定的Qdrant集合。
    注意，数据库Qdrant每个存储兑现是一个Point（点），每个点的索引内容（向量表示）如果没有上下文，那么后期搜索会很麻烦，比如这里的信息是公司的主要从事的业务，但却没有提及公司名称，
    那么后期搜索的时候，会错误的匹配，或者找不到想要的组合。所以这里使用了如下信息作为向量生成
    文档名称+文档首页（chunk）摘要+本页摘要+本页内容（1000个token）
    实际使用的效果很好，很少出现文不对题的情况。前台搜索的时候，另外加入了针对于文件名的筛选，这样就可以和某一个或者某多个文档进行对话，提问的问题就可以更宽泛（匹配最低值降低），回答的问题也可以更加准确。

6. 更新状态文件以记录处理过程。

7. 如果所有文件都已处理完毕，输出完成信息。
"""
def main_loop(collection_name, base_dir):
    # 清除掉重复的文件
    print("start to delete duplicate files")
    src_folder = base_dir  # 替换为需要遍历的文件夹
    duplicate_folder = "duplicate_delete"
    find_duplicate_files_and_move(src_folder, duplicate_folder)
    
    # connect to database, create if not exists
    client = QdrantClient("127.0.0.1", port=6333)
    try:
        # Get information about existing collection
        collection_exists = client.get_collection(collection_name)
        print(f"{collection_exists} exists")
    except Exception as e:
        print(f"Collection {collection_name} not exists, create it")
        client.recreate_collection(
            collection_name=collection_name,
            vectors_config=VectorParams(size=768, distance=Distance.COSINE),
        )
        
    status_file_path = os.path.join(base_dir, 'file_processing_status.json')
    # Initialize the status file with an empty dictionary if it does not exist
    if not os.path.exists(status_file_path):
        with open(status_file_path, 'w') as f:
            json.dump({}, f)
    
    for root, dirs, files in os.walk(base_dir):
        for file in tqdm.tqdm(files):
            file_path = os.path.join(root, file)
            file_dir = os.path.dirname(file_path)
            file_name = os.path.basename(file_path)
            if file_path == status_file_path:
                continue
            file_update_time = os.path.getmtime(file_path)
            file_uuid = str(uuid.uuid5(namespace, f"{file_path}_0")) # check file chunk 0 if exists
            
            file_status = get_file_processing_status(status_file_path, file_path)
            if file_status and file_status['is_completed']:
                print(f"Skipping already processed file: {file_path}")
                continue

            last_processed_chunk = file_status['last_processed_chunk'] if file_status else -1
            if last_processed_chunk != -1:
                print(f"Resuming processing of file: {file_path} from chunk {last_processed_chunk + 1}")

            print(f"processing file_path:{file_path} with UUiD:{file_uuid} to collection:{collection_name}")
            
            
            """
            按照如下方法进行索引：
            文件的文本分为1000个字符的chunk， 由多个chunk组成
            每个chunk的索引由一下内容组成：
            1. 文件名
            2. 第一页的chunk的summary
            3. 本页的summary
            4. 本页的内容
            然后把第一页也就是第一个chunk做summary，这个summary+文件名和其他所有的chunk每个都结合做索引，但是存储内容只有chunk自身的内容
            """
            # print(f"existing_points:{file_path}")
            # file_type = "pdf" # 将来打算再加上text, web, userinput等类型
            # if not existing_points:
                
            if file.lower().endswith('.pdf'):
                file_type = "pdf"
                chunks = get_pdf_chunks(file_path)
            elif file.lower().endswith('.docx'):
                file_type = "word"
                doc = docx.Document(file_path)
                text = '\n'.join([para.text for para in doc.paragraphs])
                chunks = split_text_to_chunks(text)
            elif file.lower().endswith('.doc'):
                print("cannot process .doc file, please convert it to .docx first, next...")
                continue
            elif file.lower().endswith('.pptx'):
                file_type = "pptx"
                text = read_pptx_text(file_path)
                chunks = split_text_to_chunks(text)
            elif file.lower().endswith('.txt'):
                file_type = "text"
                with open(file_path, 'r', encoding='utf-8') as f:
                    text = f.read()
                    chunks = split_text_to_chunks(text)
            else:
                continue
            
            chunk0_summary = None
            for file_chunk, chunk_text in enumerate(chunks):
                is_last_chunk = file_chunk == len(chunks) - 1
                point_id = str(uuid.uuid5(namespace, f"{file_path}_{file_chunk}"))
                existing_points = client.retrieve(
                    collection_name=collection_name,
                    ids=[point_id], # check current point if exists
                    with_payload=True
                )
                if existing_points:
                    # 下面这段运行一整个循环以后，就可以去掉了，因为之前索引的没有这个file_dir
                    if 'file_dir' not in existing_points[0].payload:
                        # existing_points[0].payload['file_dir'] = file_dir
                        print(f"updating...... point_id:{point_id} file_name:{file_name} file_chunk:{file_chunk} file_dir:{file_dir}")
                        client.set_payload(
                            collection_name=collection_name,
                            wait=True,
                            payload={
                                "file_dir": file_dir
                            },
                            points=[point_id]
                        )
                    if  file_chunk > last_processed_chunk:
                        # 存在此数据， 但是却没有更新到
                        update_processing_status_file(status_file_path, file_path, file_chunk, is_last_chunk)
                        last_processed_chunk = file_chunk
                else:
                    # 数据库并未存在此记录
                    if chunk0_summary is None:
                        chunk0_summary = get_summary(chunks[0])
                    chunk_summary = get_summary(chunk_text) # summary cost too much time
                    embedding_text = f"{file_name} {chunk0_summary} {chunk_summary} {chunk_text}"
                    print("---------------------embedding_text---------------------")
                    print(embedding_text)
                    embedding_vector = to_embeddings(embedding_text)
                    print(f"inserting...... point_id:{point_id} file_name:{file_name}\nchunk_summary:{chunk_summary}\nchunk_text:{chunk_text}\n ")
                    client.upsert(
                        collection_name=collection_name,
                        wait=True,
                        points=[
                            # user file_path as id, for future search
                            PointStruct(id=point_id, vector=embedding_vector, payload={
                                            "title": chunk_summary, 
                                            "text": chunk_text, 
                                            "file_summary": chunk0_summary,
                                            "file_type":file_type, 
                                            "file_update_time": file_update_time, 
                                            "file_name": file_name,
                                            "file_dir": file_dir,
                                            "file_uuid":file_uuid, 
                                            "file_chunk": file_chunk
                                            }),
                        ],
                    )
                    
                    update_processing_status_file(status_file_path, file_path, file_chunk, is_last_chunk)



parser = argparse.ArgumentParser(description="Update collections with specified collection name")
parser.add_argument("--collection_name", type=str, required=True, help="Specify the collection name")
args = parser.parse_args()
collection_name = args.collection_name

if __name__ == '__main__':
    with open('../config.json', 'r') as f:
        config = json.load(f)
        
    if collection_name not in config['base_dir']:
        raise KeyError(f"Collection '{collection_name}' is not defined in config file.")
    base_dir = config['base_dir'][collection_name]
    if not os.path.isdir(base_dir):
        print(f"Error: The base directory {base_dir} for collection '{args.collection_name}' does not exist.")
        sys.exit(1)
    print(f"process from dir {base_dir} to collection {collection_name}")
    while True:
        main_loop(collection_name, base_dir)
        time.sleep(60)
